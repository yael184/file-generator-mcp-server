"""PowerPoint file generation service."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from models.file_request import FileRequest
from models.file_response import build_file_response
from services.file_service import get_output_path
from validation import validate_filename, validate_content


def generate_pptx(request: FileRequest) -> dict:
    """
    Generate a PowerPoint presentation from content.

    Args:
        request: FileRequest with filename and content

    Returns:
        Dictionary with file metadata and commands
    """
    validate_filename(request.filename)
    validate_content(request.content)

    output_path = get_output_path(request.filename, "pptx")

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout

    slide = prs.slides.add_slide(blank_slide_layout)
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = request.filename
    p.font.size = Pt(44)
    p.font.bold = True

    # Add content on second slide
    slide2 = prs.slides.add_slide(blank_slide_layout)
    txBox2 = slide2.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = request.content
    p2.font.size = Pt(18)

    prs.save(str(output_path))

    return build_file_response(output_path)
